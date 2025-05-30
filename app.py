import streamlit as st
import tempfile, zipfile, os
import cv2
from PIL import Image, ImageDraw
from extraction import get_video_dimensions, extract_frames_with_changes
from extraction import (
    get_video_dimensions,
    extract_frames_with_changes,
    get_sample_frame,           
)
import shutil
from pptx import Presentation
from pptx.util import Emu
from ocr import ocr_image_to_dataframe

st.set_page_config(page_title="Video → PPT-Screenshots", layout="centered")
# Bild ohne abgerundete Ecken:
st.markdown(
    """
    <style>
      /* Entfernt Abrundung bei allen Bildern */
      .stImage img {
        border-radius: 0 !important;
      }
    </style>
    """,
    unsafe_allow_html=True
)

st.title("📽️ PowerPoint-Screenshots aus Video extrahieren")

# Video-Upload
video_file = st.file_uploader("Video hochladen (MP4, AVI…)", type=["mp4","avi","mov"])
if not video_file:
    st.info("Bitte lade ein Video hoch, um zu starten.")
    st.stop()

# Speichere Upload temporär
tfile = tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(video_file.name)[1])
tfile.write(video_file.read())
tfile.flush()

# --- Sample Screenshot ---
sample = get_sample_frame(tfile.name)  
sample_rgb = cv2.cvtColor(sample, cv2.COLOR_BGR2RGB)
sample_pil = Image.fromarray(sample_rgb)



# Video-Dimensionen ermitteln und ROI-Eingabe
height, width = get_video_dimensions(tfile.name)
st.markdown(f"**Video-Größe:** {width}×{height} Pixel")


col1, col2 = st.columns(2)
with col1:
    x = st.number_input(
        "Startpunkt Frame-Breite in Pixel",
        min_value=0,
        max_value=width - 1,
        value=0
    )
    max_w = width - x
    w = st.number_input(
        "Schrittweite Frame-Breite in Pixel",
        min_value=1,
        max_value=max_w,
        value=max_w        # <— Default jetzt dynamisch
    )

with col2:
    y = st.number_input(
        "Startpunkt Frame-Weite in Pixel",
        min_value=0,
        max_value=height - 1,
        value=0
    )
    max_h = height - y
    h = st.number_input(
        "Schrittweite Frame-Breite in Pixel",
        min_value=1,
        max_value=max_h,
        value=max_h        # <— Default dynamisch
    )



# Sobald ROI definiert ist, zeichne sie in die Vorschau
overlay = sample_pil.copy()
draw    = ImageDraw.Draw(overlay)
# Rechteck um die ROI
draw.rectangle([x, y, x+w, y+h], outline="red", width=3)
# Gitter-Linien
draw.line([(x,   0),      (x,   height)],   fill="red", width=1)
draw.line([(x+w, 0),      (x+w, height)],   fill="red", width=1)
draw.line([(0,   y),      (width, y)],      fill="red", width=1)
draw.line([(0,   y+h),    (width, y+h)],    fill="red", width=1)

# Dicker Punkt bei (x, y) als Referenz-Anker
r = 8  # Radius des Punktes
draw.ellipse(
    [(x - r, y - r), (x + r, y + r)],
    fill="red",
    outline="red"
)

st.subheader("🔲 Vorschau Beispielbild mit Frame")
st.image(overlay, use_container_width =True)

threshold = st.slider("Änderungssensitivität (0–1)", min_value=0.0, max_value=1.0, value=0.1, step=0.01)
interval  = st.number_input("Prüf-Intervall (Sekunden)", value=5, min_value=1)


# Init session state
if "outs" not in st.session_state:
    st.session_state["outs"] = None
if "ocr_ready" not in st.session_state:
    st.session_state["ocr_ready"] = False

# … Dein Upload-, ROI- und Extraktions-Block oben …
# 1) Extraktion-Button
if st.button("Screenshots extrahieren"):
    with st.spinner("Extrahiere…"):
        roi = (y, x, h, w)
        outs = extract_frames_with_changes(
            tfile.name, roi, threshold, check_interval_s=interval
        )

    if not outs:
        st.warning("Keine signifikanten Änderungen gefunden.")
    else:
        st.success(f"{len(outs)} Screenshots erstellt.")
        # Zustand speichern
        st.session_state["outs"] = outs
        st.session_state["ocr_ready"] = True

        # ZIP mit Screenshots anbieten
        zip_path = os.path.join(os.path.dirname(tfile.name), "screenshots.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            for fp in outs:
                zf.write(fp, arcname=os.path.basename(fp))
        with open(zip_path, "rb") as f:
            st.download_button(
                "📥 Screenshots (ZIP) herunterladen",
                f,
                file_name="screenshots.zip"
            )

if st.session_state["ocr_ready"]:
    if st.button("📝 OCR ausführen & PPTX generieren"):
        with st.spinner("OCR & PPTX wird erstellt…"):
            outs = sorted(st.session_state["outs"])  # zeitlich sortieren
            base_dir = os.path.dirname(tfile.name)
            ocr_dir  = os.path.join(base_dir, "ocr_output")
            os.makedirs(ocr_dir, exist_ok=True)

            # 1) Neue Präsentation anlegen
            prs = Presentation()
            emu_x = emu_y = None

            for fp in outs:
                # --- a) Bild-Folie ---
                slide_img = prs.slides.add_slide(prs.slide_layouts[6])
                # Vollbild einfügen
                slide_img.shapes.add_picture(
                    fp,
                    left=0, top=0,
                    width=prs.slide_width,
                    height=prs.slide_height
                )

                # EMU-Faktor fürs spätere Text-Layout nur einmal bestimmen
                if emu_x is None:
                    img = cv2.imread(fp)
                    h, w = img.shape[:2]
                    emu_x = prs.slide_width  / w
                    emu_y = prs.slide_height / h

                # --- b) OCR-Text-Folie ---
                df_ocr = ocr_image_to_dataframe(fp)
                # Textblöcke extrahieren
                blocks = []
                for _, grp in df_ocr.groupby("block_num"):
                    g = grp.sort_values(["line_num", "word_num"])
                    text = " ".join(g["text"].tolist())
                    left   = int(g["left"].min())
                    top    = int(g["top"].min())
                    right  = int((g["left"] + g["width"]).max())
                    bottom = int((g["top"]  + g["height"]).max())
                    blocks.append({
                        "text":   text,
                        "left":   left,
                        "top":    top,
                        "width":  right - left,
                        "height": bottom - top
                    })

                slide_txt = prs.slides.add_slide(prs.slide_layouts[6])
                for b in blocks:
                    tb = slide_txt.shapes.add_textbox(
                        Emu(b["left"]  * emu_x),
                        Emu(b["top"]   * emu_y),
                        Emu(b["width"] * emu_x),
                        Emu(b["height"]* emu_y),
                    )
                    tf = tb.text_frame
                    tf.text = b["text"]
                    for p in tf.paragraphs:
                        p.font.size = Emu(12 * 9144)

            # 2) PPTX speichern
            pptx_path = os.path.join(ocr_dir, "recreated_presentation.pptx")
            prs.save(pptx_path)

            # 3) ZIP mit allen PNGs + PPTX
            zip2 = os.path.join(base_dir, "ocr_slides.zip")
            with zipfile.ZipFile(zip2, "w") as zf2:
                # alle Screenshots
                for fp in outs:
                    zf2.write(fp, arcname=os.path.basename(fp))
                # und die Präsentation
                zf2.write(pptx_path, arcname=os.path.basename(pptx_path))

        # 4) Download-Button
        with open(zip2, "rb") as f2:
            st.download_button(
                "📥 Screenshots + OCR-PPTX herunterladen",
                f2,
                file_name="export_slides.zip"
            )
