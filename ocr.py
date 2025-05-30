import os
import cv2
import pandas as pd
import pytesseract
from pytesseract import Output
from PIL import Image, ImageDraw
import shutil
from pptx import Presentation
from pptx.util import Emu

def ocr_image_to_dataframe(image_path: str) -> pd.DataFrame:
    """
    Lädt das Bild unter `image_path`, führt OCR mit pytesseract durch
    und gibt ein DataFrame zurück mit den erkannten Textblöcken und Metadaten:
    block_num, par_num, line_num, word_num, left, top, width, height, conf, text.
    """
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"Image not found: {image_path}")

    # Bild laden (BGR) und prüfen
    img_bgr = cv2.imread(image_path)
    if img_bgr is None:
        raise IOError(f"Failed to read image: {image_path}")

    # In RGB für pytesseract konvertieren
    img_rgb = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2RGB)

    # OCR durchführen
    data = pytesseract.image_to_data(img_rgb, output_type=Output.DICT)

    # Nur Einträge mit nicht-leerem Text übernehmen
    rows = []
    n_boxes = len(data["level"])
    for i in range(n_boxes):
        text = data["text"][i].strip()
        if text:
            rows.append({
                "block_num": data["block_num"][i],
                "par_num":   data["par_num"][i],
                "line_num":  data["line_num"][i],
                "word_num":  data["word_num"][i],
                "left":      data["left"][i],
                "top":       data["top"][i],
                "width":     data["width"][i],
                "height":    data["height"][i],
                "conf":      data["conf"][i],
                "text":      text
            })

    return pd.DataFrame(rows)

def add_image_and_text_slides(prs: Presentation, image_path: str):
    """
    Fügt dem Presentation-Objekt zwei Folien hinzu:
      1) Vollbild mit dem Screenshot
      2) Bearbeitbare Textfolie basierend auf OCR-Ergebnissen
    """
    # Bild-Folie
    slide_img = prs.slides.add_slide(prs.slide_layouts[6])
    slide_img.shapes.add_picture(
        image_path,
        left=0, top=0,
        width=prs.slide_width,
        height=prs.slide_height
    )

    # EMU-Skala bestimmen
    img = cv2.imread(image_path)
    h, w = img.shape[:2]
    emu_x = prs.slide_width  / w
    emu_y = prs.slide_height / h

    # OCR-Text-Folie
    df = ocr_image_to_dataframe(image_path)
    slide_txt = prs.slides.add_slide(prs.slide_layouts[6])
    for _, grp in df.groupby("block_num"):
        g = grp.sort_values(["line_num", "word_num"])
        text = " ".join(g["text"].tolist())
        left   = int(g["left"].min())
        top    = int(g["top"].min())
        right  = int((g["left"] + g["width"]).max())
        bottom = int((g["top"]  + g["height"]).max())

        tb = slide_txt.shapes.add_textbox(
            Emu(left * emu_x),
            Emu(top  * emu_y),
            Emu((right-left) * emu_x),
            Emu((bottom-top) * emu_y),
        )
        tf = tb.text_frame
        tf.text = text
        for p in tf.paragraphs:
            p.font.size = Emu(12 * 9144)